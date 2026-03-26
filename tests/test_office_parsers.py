"""Tests for office-document parser support in ingestion."""

from __future__ import annotations

from pathlib import Path

from coderag.ingestion.repo_scanner import scan_folder
from coderag.parsers.generic_parser import parse_by_extension


def test_scan_folder_includes_new_office_extensions(tmp_path: Path) -> None:
    """Ensure scanner detects doc, pptx and xlsx files."""
    (tmp_path / "legacy.doc").write_bytes(b"Simple DOC text")
    (tmp_path / "slides.pptx").write_bytes(b"placeholder")
    (tmp_path / "sheet.xlsx").write_bytes(b"placeholder")

    files = scan_folder(tmp_path)
    names = {path.name for path in files}

    assert "legacy.doc" in names
    assert "slides.pptx" in names
    assert "sheet.xlsx" in names


def test_parse_doc_returns_text_for_binary_like_input(tmp_path: Path) -> None:
    """DOC parser should return best-effort text from binary content."""
    file_path = tmp_path / "legacy.doc"
    file_path.write_bytes(b"\x00\x01This is a legacy doc content\x02\x03")

    parsed = parse_by_extension(file_path)

    assert isinstance(parsed, str)
    assert "legacy doc" in parsed.lower()


def test_parse_pptx_extracts_slide_text(tmp_path: Path) -> None:
    """PPTX parser should extract visible text from slides."""
    from pptx import Presentation

    file_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Architecture Overview"
    slide.placeholders[1].text = "Hybrid retrieval and graph expansion"
    presentation.save(str(file_path))

    parsed = parse_by_extension(file_path)

    assert "Architecture Overview" in parsed
    assert "Hybrid retrieval" in parsed


def test_parse_xlsx_extracts_cell_values(tmp_path: Path) -> None:
    """XLSX parser should extract non-empty worksheet cells."""
    from openpyxl import Workbook

    file_path = tmp_path / "book.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet["A1"] = "Project"
    sheet["B1"] = "Amount"
    sheet["A2"] = "Atlas"
    sheet["B2"] = 1200
    workbook.save(str(file_path))

    parsed = parse_by_extension(file_path)

    assert "[Sheet] Budget" in parsed
    assert "Project" in parsed
    assert "Atlas" in parsed
    assert "1200" in parsed
