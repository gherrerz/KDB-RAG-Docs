"""PPTX parser for extracting slide text."""

from __future__ import annotations

from pathlib import Path


def parse_pptx(file_path: Path) -> str:
    """Extract text content from a PPTX presentation.

    Returns an empty string when parsing fails so ingestion can continue.
    """
    try:
        from pptx import Presentation

        presentation = Presentation(str(file_path))
        lines: list[str] = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        lines.append(text)
                table = getattr(shape, "table", None)
                if table is not None:
                    for row in table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                lines.append(text)

        return "\n".join(lines)
    except Exception:
        return ""
